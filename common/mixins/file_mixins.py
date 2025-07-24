from io import StringIO
import os
from rest_framework import serializers
from pprint import pprint
from dotenv import load_dotenv
from pptx import Presentation
from docx import Document

load_dotenv()


class FileValidationMixin:
    def validate_file(self, file, max_size=int(os.getenv("MAX_FILE_SIZE_MB")) * 1024 * 1024, allowed_mime_types=None):
        if file.size > max_size:
            raise serializers.ValidationError(f"{file.name} cannot exceed size limit of {max_size} MB")
    
        if allowed_mime_types is None:
            allowed_mime_types = [
                "application/msword",  # .doc
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # .docx
                "application/vnd.ms-powerpoint",  # .ppt
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
                "application/pdf",  # .pdf
            ]
       
        if file.content_type not in allowed_mime_types:
            raise serializers.ValidationError(f"{file.name} has an unsupported type of {file.content_type}")

        return file


class ExtractContent:
    def extract_document(content_type, document):
        match content_type:
            case ("application/msword" | "application/vnd.openxmlformats-officedocument.wordprocessingml.document"):
                return ExtractContent.extract_docx(document)
            case "application/pdf":
                return ""
            case 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                return ExtractContent.extract_pptx(document)
            
    def extract_docx(file):
        doc = Document(file)
        text = []

        for para in doc.paragraphs:
            if para.text.strip():  # Skip empty lines
                text.append(para.text.strip())

        return '\n'.join(text)
    
    def extract_pptx(file):
        prs = Presentation(file)
        text_list = []

        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text.strip())
            text_list.append("".join(slide_text))
        
        return '\n'.join(text_list)